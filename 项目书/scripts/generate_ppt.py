"""
Generate completion report PPT for course project.
Based on: 项目书/最终汇报_实验分析与结论.md
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Color Palette ──
NAVY = RGBColor(0x1A, 0x3A, 0x5C)
TEAL = RGBColor(0x2E, 0x86, 0xAB)
AMBER = RGBColor(0xF1, 0x8F, 0x01)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE = RGBColor(0xF5, 0xF7, 0xFA)
DARK = RGBColor(0x2C, 0x3E, 0x50)
GRAY = RGBColor(0x7F, 0x8C, 0x9A)
LIGHT_GRAY = RGBColor(0xE8, 0xEC, 0xF0)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xE7, 0x4C, 0x3C)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        import lxml.etree as etree
        # alpha not trivial in python-pptx; skip
    return shape

def add_text_box(slide, left, top, w, h, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_para(text_frame, text, font_size=16, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri", space_before=None, space_after=None):
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if space_before:
        p.space_before = Pt(space_before)
    if space_after:
        p.space_after = Pt(space_after)
    return p

def add_rich_para(text_frame, segments, alignment=PP_ALIGN.LEFT, space_before=None):
    """segments: list of (text, font_size, color, bold)"""
    p = text_frame.add_paragraph()
    p.alignment = alignment
    if space_before:
        p.space_before = Pt(space_before)
    for seg in segments:
        run = p.add_run()
        run.text = seg[0]
        run.font.size = Pt(seg[1])
        run.font.color.rgb = seg[2]
        run.font.bold = seg[3] if len(seg) > 3 else False
        run.font.name = "Calibri"
    return p

def make_slide():
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    return slide

def slide_number(slide, num, total):
    add_text_box(slide, W - Inches(1.2), H - Inches(0.5), Inches(1), Inches(0.4),
                 f"{num}/{total}", 10, GRAY, alignment=PP_ALIGN.RIGHT)

TOTAL_SLIDES = 18

# ══════════════════════════════════════════════════════
# SLIDE 1: Title
# ══════════════════════════════════════════════════════
slide = make_slide()
add_bg(slide, NAVY)
add_shape(slide, 0, Inches(2.8), W, Inches(3.5), TEAL)
add_text_box(slide, Inches(1), Inches(1.0), Inches(11), Inches(0.6),
             "课程作业 · 软件体系结构", 16, WHITE, bold=False)
add_text_box(slide, Inches(1), Inches(3.1), Inches(11), Inches(1.2),
             "神经网络鲁棒性验证策略对比", 40, WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(4.3), Inches(11), Inches(0.8),
             "α,β-CROWN 与 Marabou 在 MNIST/CIFAR-10 上的实验分析", 20, OFF_WHITE)
add_text_box(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.5),
             "完成情况汇报 · 2026-05", 14, RGBColor(0xBB, 0xCC, 0xDD))
slide_number(slide, 1, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════
# SLIDE 2: Outline
# ══════════════════════════════════════════════════════
slide = make_slide()
add_bg(slide, WHITE)
add_shape(slide, 0, 0, Inches(0.15), H, TEAL)
add_text_box(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.7),
             "汇报提纲", 32, NAVY, bold=True)

items = [
    ("1", "问题定义与实验框架", "验证问题、指标、模型与数据集"),
    ("2", "轻量方法：能力边界", "CROWN / α-CROWN / PGD — 原理、实验与局限"),
    ("3", "完整验证：策略优化", "BaB 三策略对比 → ADO 方法论 → 跨 ε 稳定性"),
    ("4", "外部验证", "CIFAR-10 跨数据集 · Marabou 跨工具"),
    ("5", "论文复现对照", "α,β-CROWN ICLR'21 + VNN-COMP'21 — 5/5 核心结论"),
    ("6", "综合结论与工程推荐", "能力光谱 · ε 分区方案"),
]
for i, (num, title, desc) in enumerate(items):
    y = Inches(1.6) + Inches(i * 0.85)
    add_shape(slide, Inches(0.8), y, Inches(0.6), Inches(0.6), TEAL)
    add_text_box(slide, Inches(0.85), y + Inches(0.08), Inches(0.5), Inches(0.5),
                 num, 22, WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1.7), y + Inches(0.02), Inches(8), Inches(0.4),
                 title, 20, NAVY, bold=True)
    add_text_box(slide, Inches(1.7), y + Inches(0.38), Inches(8), Inches(0.3),
                 desc, 14, GRAY)
slide_number(slide, 2, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════
# SLIDE 3: Problem Definition
# ══════════════════════════════════════════════════════
slide = make_slide()
add_bg(slide, OFF_WHITE)
add_shape(slide, 0, 0, W, Inches(1.0), NAVY)
add_text_box(slide, Inches(0.8), Inches(0.15), Inches(10), Inches(0.7),
             "问题定义与实验框架", 28, WHITE, bold=True)

# Left: problem definition
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5.5), Inches(0.4),
             "核心验证问题", 20, TEAL, bold=True)
tb = add_text_box(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.5),
             "给定输入 x、真实标签 y、扰动半径 ε (L∞)，判定：", 14, DARK)
tf = tb.text_frame
tf.word_wrap = True
add_para(tf, "∀x' : ‖x'−x‖∞ ≤ ε, x'∈[0,1] → f_y(x') > f_j(x'), ∀j≠y", 13, NAVY, bold=True)
add_para(tf, "若成立 → safe（局部鲁棒）；存在反例 → unsafe；超时 → unknown", 13, DARK)

add_text_box(slide, Inches(0.8), Inches(4.5), Inches(5.5), Inches(0.4),
             "评价指标", 20, TEAL, bold=True)
tb2 = add_text_box(slide, Inches(0.8), Inches(5.0), Inches(5.5), Inches(2.0),
             "", 14, DARK)
tf2 = tb2.text_frame
tf2.word_wrap = True
add_para(tf2, "• VRA (Verified Accuracy) = Nsafe / Ntotal — 核心质量指标", 14, DARK)
add_para(tf2, "• Timeout 数 — 极端样本占比", 14, DARK)
add_para(tf2, "• Mean Time — 整体计算开销", 14, DARK)

# Right: experimental object box
add_shape(slide, Inches(7.0), Inches(1.3), Inches(5.5), Inches(5.5), WHITE)
add_text_box(slide, Inches(7.4), Inches(1.5), Inches(4.8), Inches(0.4),
             "实验对象", 20, TEAL, bold=True)

info_items = [
    ("模型架构", "FCNN: 784→256→128→10, ReLU ×2"),
    ("参数量", "~260K"),
    ("训练集", "MNIST: 60000/10000, Clean Acc ≈99%"),
    ("扰动范数", "L∞，各 ε: 0.01 / 0.02 / 0.03 / 0.05"),
    ("模型格式", "ONNX（两工具原生读取）"),
    ("数据集", "MNIST（主线）+ CIFAR-10（迁移验证）"),
    ("验证工具", "α,β-CROWN（主）+ Marabou SMT（对照）"),
]
for i, (k, v) in enumerate(info_items):
    y = Inches(2.2) + Inches(i * 0.65)
    add_text_box(slide, Inches(7.4), y, Inches(1.5), Inches(0.4),
                 k, 13, GRAY, bold=False)
    add_text_box(slide, Inches(9.0), y, Inches(3.2), Inches(0.4),
                 v, 13, DARK, bold=False)

slide_number(slide, 3, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════
# SLIDE 4: Lightweight Methods Overview
# ══════════════════════════════════════════════════════
slide = make_slide()
add_bg(slide, OFF_WHITE)
add_shape(slide, 0, 0, W, Inches(1.0), NAVY)
add_text_box(slide, Inches(0.8), Inches(0.15), Inches(10), Inches(0.7),
             "轻量方法层：边界传播与经验攻击", 28, WHITE, bold=True)

# Three method cards
methods = [
    ("CROWN", "固定线性松弛 + 反向传播",
     [("VRA ε=0.01→0.05", "98% → 4%"),
      ("耗时", "~0.2s 恒定，与 ε 无关"),
      ("特性", "极快但不完备")],
     NAVY),
    ("α-CROWN", "可优化线性松弛 (α参数)",
     [("VRA 提升", "比 CROWN +1~5pp"),
      ("最佳收益点", "ε=0.03 (+5pp)"),
      ("特性", "倒U型收益曲线")],
     TEAL),
    ("PGD", "投影梯度下降 ×30重启",
     [("Unsafe 检出", "0% → 58% (ε↑)"),
      ("召回率", "覆盖 BaB 全部反例"),
      ("特性", "能找反例，不能证安全")],
     AMBER),
]
for i, (name, subtitle, facts, color) in enumerate(methods):
    x = Inches(0.6) + Inches(i * 4.2)
    add_shape(slide, x, Inches(1.4), Inches(3.8), Inches(5.5), WHITE)
    add_shape(slide, x, Inches(1.4), Inches(3.8), Inches(0.8), color)
    add_text_box(slide, x + Inches(0.3), Inches(1.5), Inches(3.2), Inches(0.5),
                 name, 26, WHITE, bold=True)
    add_text_box(slide, x + Inches(0.3), Inches(2.0), Inches(3.2), Inches(0.4),
                 subtitle, 12, OFF_WHITE)
    for j, (k, v) in enumerate(facts):
        y = Inches(2.6) + Inches(j * 1.2)
        add_text_box(slide, x + Inches(0.3), y, Inches(3.2), Inches(0.3),
                     k, 12, GRAY)
        add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(3.2), Inches(0.5),
                     v, 18, DARK, bold=True)

slide_number(slide, 4, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════
# SLIDE 5: BaB + Branching Strategies
# ══════════════════════════════════════════════════════
slide = make_slide()
add_bg(slide, OFF_WHITE)
add_shape(slide, 0, 0, W, Inches(1.0), NAVY)
add_text_box(slide, Inches(0.8), Inches(0.15), Inches(10), Inches(0.7),
             "完整验证：BaB 分支定界与三策略对比", 28, WHITE, bold=True)

# Strategy comparison table
table_data = [
    ["策略", "VRA", "safe", "timeout", "Mean Time", "Max Time"],
    ["baseline (babsr)", "91.0%", "91", "9", "3.82s", "49.17s"],
    ["auto", "91.0%", "91", "9", "5.77s", "51.56s"],
    ["kfsb", "92.0%", "92", "8", "3.17s", "34.27s"],
]
rows, cols = len(table_data), len(table_data[0])
tbl = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.3), Inches(11.5), Inches(2.5)).table
for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.text = table_data[r][c]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.size = Pt(13)
            paragraph.font.name = "Calibri"
            if r == 0:
                paragraph.font.bold = True
                paragraph.font.color.rgb = WHITE
            elif r == 3:
                paragraph.font.bold = True
                paragraph.font.color.rgb = NAVY
            else:
                paragraph.font.color.rgb = DARK
        # cell color
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
        elif r == 3:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xF8)
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE

# Key insight box
add_shape(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(2.8), WHITE)
add_text_box(slide, Inches(1.2), Inches(4.4), Inches(10.8), Inches(0.4),
             "关键发现", 20, TEAL, bold=True)

bullets = [
    "kfsb 的 lookahead 评分机制全面优于 babsr 的局部偏差评分：VRA +1pp, timeout −1, Mean Time −0.65s",
    "auto 策略 VRA 与 baseline 持平但耗时高 51% — 自适应决策 overhead 未被精度提升补偿",
    "kfsb 的 Max Time 降低 30% — 对最难样本的处理能力显著增强",
]
tb = add_text_box(slide, Inches(1.2), Inches(5.0), Inches(10.8), Inches(1.8), "", 14, DARK)
tf = tb.text_frame
tf.word_wrap = True
for b in bullets:
    add_para(tf, "▸  " + b, 14, DARK, space_before=6)

slide_number(slide, 5, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════
# SLIDE 6: ADO Methodology (KEY SLIDE)
# ══════════════════════════════════════════════════════
slide = make_slide()
add_bg(slide, NAVY)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "核心改进：受控消融驱动的验证工具优化方法论 (ADO)", 28, WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.4),
             "Ablation-Driven Optimization — 从"调参"到"系统性分析方法论"", 14, RGBColor(0xBB, 0xCC, 0xDD))

# Three steps
steps = [
    ("1", "维度分解", "识别 3 个可调参数\nmethod × candidates × reduceop\n建立消融表"),
    ("2", "受控测量", "单变量法固定其他维度\n测量目标参数的独立效应\nΔVRA / ΔTimeout / ΔMeanTime"),
    ("3", "量化决策", "根据边际贡献组合最优\nmethod=kfsb, candidates=5\nreduceop=min"),
]
for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.6) + Inches(i * 4.2)
    add_shape(slide, x, Inches(1.7), Inches(3.8), Inches(3.0), TEAL)
    add_shape(slide, x + Inches(1.4), Inches(1.3), Inches(1.0), Inches(0.8), AMBER)
    add_text_box(slide, x + Inches(1.1), Inches(1.35), Inches(1.6), Inches(0.7),
                 num, 32, WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), Inches(2.4), Inches(3.2), Inches(0.4),
                 title, 20, WHITE, bold=True)
    add_text_box(slide, x + Inches(0.3), Inches(2.9), Inches(3.2), Inches(1.5),
                 desc, 13, OFF_WHITE)

# Ablation decomposition table
add_shape(slide, Inches(0.6), Inches(5.0), Inches(12.0), Inches(2.2), RGBColor(0x15, 0x30, 0x4F))
add_text_box(slide, Inches(1.0), Inches(5.1), Inches(5), Inches(0.3),
             "消融效应分解", 16, AMBER, bold=True)

ablation_data = [
    ["参数变动", "对照对", "Δ VRA", "Δ Timeout", "Δ Mean Time", "效应评估"],
    ["method: babsr → kfsb", "baseline vs kfsb", "+1.0pp", "−1", "−0.46s", "主导因子"],
    ["candidates: 10 → 5", "kfsb vs kfsb_c5", "+1.0pp", "−1", "−0.35s", "显著正向"],
    ["reduceop: min → max", "kfsb vs kfsb_max", "0", "0", "−0.15s", "微弱"],
]
r2, c2 = len(ablation_data), len(ablation_data[0])
tbl2 = slide.shapes.add_table(r2, c2, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.6)).table
for r in range(r2):
    for c in range(c2):
        cell = tbl2.cell(r, c)
        cell.text = ablation_data[r][c]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(12)
            p.font.name = "Calibri"
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = AMBER
            elif r == 1:
                p.font.bold = True
                p.font.color.rgb = GREEN
            elif r == 2:
                p.font.bold = True
                p.font.color.rgb = GREEN
            else:
                p.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY if r > 0 else RGBColor(0x10, 0x25, 0x40)

slide_number(slide, 6, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════
# SLIDE 7: ADO Result - Big Numbers
# ══════════════════════════════════════════════════════
slide = make_slide()
add_bg(slide, OFF_WHITE)
add_shape(slide, 0, 0, W, Inches(1.0), NAVY)
add_text_box(slide, Inches(0.8), Inches(0.15), Inches(10), Inches(0.7),
             "ADO 优化结果：kfsb_candidates5 vs baseline", 28, WHITE, bold=True)

# Big stat callouts
stats = [
    ("VRA", "+2.0pp", "93.0%", "91.0%", GREEN),
    ("Timeout", "−2", "7", "9", AMBER),
    ("Mean Time", "−0.81s", "3.24s", "4.06s", TEAL),
]
for i, (label, delta, new_val, old_val, color) in enumerate(stats):
    x = Inches(0.8) + Inches(i * 4.1)
    add_shape(slide, x, Inches(1.4), Inches(3.6), Inches(2.2), WHITE)
    add_shape(slide, x + Inches(2.6), Inches(1.0), Inches(1.0), Inches(0.8), color)
    add_text_box(slide, x + Inches(0.3), Inches(1.5), Inches(3.0), Inches(0.4),
                 label, 14, GRAY)
    add_text_box(slide, x + Inches(0.3), Inches(1.9), Inches(3.0), Inches(0.8),
                 new_val, 36, DARK, bold=True)
    add_text_box(slide, x + Inches(0.3), Inches(2.7), Inches(3.0), Inches(0.4),
                 f"baseline: {old_val}  |  {delta}", 13, color, bold=True)

# Explanation section
add_shape(slide, Inches(0.8), Inches(4.0), Inches(11.5), Inches(3.0), WHITE)
add_text_box(slide, Inches(1.2), Inches(4.2), Inches(10.8), Inches(0.4),
             "ADO vs 调参：本质区别", 20, NAVY, bold=True)

tb_explain = add_text_box(slide, Inches(1.2), Inches(4.8), Inches(4.8), Inches(2.0), "", 13, DARK)
tf = tb_explain.text_frame
tf.word_wrap = True
add_para(tf, "❌ 调参的目标：找到一组能用的数", 13, RED, bold=False, space_before=4)
add_para(tf, "  输出：一个配置点", 13, GRAY, space_before=2)
add_para(tf, "", 8, DARK)
add_para(tf, "✅ ADO 的目标：理解每个参数为什么影", 13, GREEN, bold=False, space_before=4)
add_para(tf, "   响结果、影响多大、如何交互", 13, GREEN)
add_para(tf, "  输出：可迁移的分析框架 + 可检验", 13, GRAY)
add_para(tf, "   的机制假说", 13, GRAY)

add_text_box(slide, Inches(7.0), Inches(4.8), Inches(4.8), Inches(2.0),
             "反直觉发现：candidates=5 > 10\n\nlookahead 评分存在噪声饱和\n候选集超过阈值后，额外候选的\n评分可靠性下降，稀释决策质量\n\n→ 为其他 lookahead 策略提供\n   可迁移的设计启示",
             13, DARK)

slide_number(slide, 7, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════
# SLIDE 8: Cross-ε Stability
# ══════════════════════════════════════════════════════
slide = make_slide()
add_bg(slide, OFF_WHITE)
add_shape(slide, 0, 0, W, Inches(1.0), NAVY)
add_text_box(slide, Inches(0.8), Inches(0.15), Inches(10), Inches(0.7),
             "跨扰动半径稳定性验证", 28, WHITE, bold=True)

# VRA table
vra_data = [["ε", "baseline VRA", "auto VRA", "kfsb VRA", "kfsb − baseline"],
            ["0.01", "100.0%", "100.0%", "100.0%", "0 (天花板)"],
            ["0.02", "91.0%", "91.0%", "92.0%", "+1.0pp"],
            ["0.03", "65.0%", "62.0%", "68.0%", "+3.0pp"],
            ["0.05", "12.0%", "5.0%", "11.0%", "−1.0pp"]]
r3, c3 = len(vra_data), len(vra_data[0])
tbl3 = slide.shapes.add_table(r3, c3, Inches(0.8), Inches(1.3), Inches(11.5), Inches(2.5)).table
for r in range(r3):
    for c in range(c3):
        cell = tbl3.cell(r, c)
        cell.text = vra_data[r][c]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(13)
            p.font.name = "Calibri"
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
            elif c == 4 and r == 4:
                p.font.color.rgb = RED
            elif c == 4 and r > 0:
                p.font.color.rgb = GREEN
                p.font.bold = True
            else:
                p.font.color.rgb = DARK
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY if r == 0 else WHITE

# Mean time
add_text_box(slide, Inches(0.8), Inches(4.0), Inches(5), Inches(0.4),
             "平均耗时对比", 18, TEAL, bold=True)
time_data = [["ε", "baseline", "auto", "kfsb"],
             ["0.01", "0.55s", "0.43s", "0.31s"],
             ["0.02", "3.93s", "5.36s", "3.27s"],
             ["0.03", "14.27s", "22.84s", "6.66s"],
             ["0.05*", "18.79s", "23.37s", "11.38s"]]
r4, c4 = len(time_data), len(time_data[0])
tbl4 = slide.shapes.add_table(r4, c4, Inches(0.8), Inches(4.5), Inches(5.5), Inches(2.5)).table
for r in range(r4):
    for c in range(c4):
        cell = tbl4.cell(r, c)
        cell.text = time_data[r][c]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(12)
            p.font.name = "Calibri"
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
            else:
                p.font.color.rgb = DARK
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY if r == 0 else WHITE

# Insight box
add_shape(slide, Inches(7.0), Inches(4.0), Inches(5.5), Inches(3.0), WHITE)
add_text_box(slide, Inches(7.4), Inches(4.2), Inches(4.8), Inches(0.4),
             "关键趋势", 18, TEAL, bold=True)
tb_i = add_text_box(slide, Inches(7.4), Inches(4.7), Inches(4.8), Inches(2.0), "", 14, DARK)
tf_i = tb_i.text_frame
tf_i.word_wrap = True
add_para(tf_i, "▸  kfsb 优势呈倒U型：ε=0.03 达峰值", 13, DARK, space_before=4)
add_para(tf_i, "   (VRA +3pp, 时间仅 baseline 的 47%)", 13, GRAY, space_before=2)
add_para(tf_i, "▸  ε=0.01 天花板效应，3 策略均 100%", 13, DARK, space_before=4)
add_para(tf_i, "▸  ε=0.05 超时主导，策略差异被淹没", 13, DARK, space_before=4)
add_para(tf_i, "", 8, DARK)
add_para(tf_i, "结论：分支策略优化在中等难度区间", 13, NAVY, bold=True, space_before=4)
add_para(tf_i, "最显著 — 太容易不需要，太难不够", 13, NAVY, bold=True)

slide_number(slide, 8, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════
# SLIDE 9: CIFAR-10 Cross-dataset
# ══════════════════════════════════════════════════════
slide = make_slide()
add_bg(slide, OFF_WHITE)
add_shape(slide, 0, 0, W, Inches(1.0), NAVY)
add_text_box(slide, Inches(0.8), Inches(0.15), Inches(10), Inches(0.7),
             "跨数据集验证：CIFAR-10", 28, WHITE, bold=True)

# Comparison table
cmp_data = [["指标", "MNIST (ε=0.01)", "CIFAR-10 (ε=2/255)", "差距"],
            ["CROWN VRA", "98.0%", "0.0%", "断层"],
            ["BaB VRA", "100.0%", "0.0%", "断层"],
            ["单样本耗时", "~3s", "~68s", "~23×"],
            ["PGD unsafe 检出", "0%", "80%", "—"]]
r5, c5 = len(cmp_data), len(cmp_data[0])
tbl5 = slide.shapes.add_table(r5, c5, Inches(0.8), Inches(1.3), Inches(11.5), Inches(2.5)).table
for r in range(r5):
    for c in range(c5):
        cell = tbl5.cell(r, c)
        cell.text = cmp_data[r][c]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(13)
            p.font.name = "Calibri"
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
            elif c == 3 and r in [1,2]:
                p.font.color.rgb = RED
                p.font.bold = True
            else:
                p.font.color.rgb = DARK
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY if r == 0 else WHITE

# Five-factor analysis
add_shape(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(2.8), WHITE)
add_text_box(slide, Inches(1.2), Inches(4.4), Inches(6), Inches(0.4),
             "难度骤降的 5 因素复合作用", 20, TEAL, bold=True)

factors = [
    ("输入维度", "784 → 3072 (×3.9)"),
    ("RGB 通道", "灰度 → 3 通道扰动组合爆炸"),
    ("Conv 松弛放大", "共享感受野使不确定性相互放大"),
    ("ReLU 层数", "2 → 4，累积误差非线性增长"),
    ("数据复杂度", "手写数字 → 自然图像高纹理"),
]
for i, (k, v) in enumerate(factors):
    y = Inches(4.9) + Inches(i * 0.42)
    add_text_box(slide, Inches(1.2), y, Inches(2.0), Inches(0.3),
                 f"  {i+1}. {k}", 12, NAVY, bold=True)
    add_text_box(slide, Inches(3.2), y, Inches(8.5), Inches(0.3),
                 v, 12, DARK)

slide_number(slide, 9, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════
# SLIDE 10: Marabou Cross-tool
# ══════════════════════════════════════════════════════
slide = make_slide()
add_bg(slide, OFF_WHITE)
add_shape(slide, 0, 0, W, Inches(1.0), NAVY)
add_text_box(slide, Inches(0.8), Inches(0.15), Inches(10), Inches(0.7),
             "跨工具验证：α,β-CROWN vs Marabou SMT", 28, WHITE, bold=True)

# Results summary
add_text_box(slide, Inches(0.8), Inches(1.3), Inches(5), Inches(0.4),
             "三扰动半径对照结果", 20, TEAL, bold=True)

summary_data = [
    ["ε", "一致率", "Marabou avg", "α,β-CROWN avg", "关键发现"],
    ["0.01", "5/5 ✅", "4.46s", "0.50s", "完美一致，低扰动"],
    ["0.02", "4/5 ⚠️", "3.53s", "0.40s", "首次分歧：Marabou 证出"],
    ["0.03", "4/5 ⚠️", "45.73s*", "0.41s", "分歧扩大 + 共同边界"],
]
r6, c6 = len(summary_data), len(summary_data[0])
tbl6 = slide.shapes.add_table(r6, c6, Inches(0.8), Inches(1.8), Inches(11.5), Inches(2.2)).table
for r in range(r6):
    for c in range(c6):
        cell = tbl6.cell(r, c)
        cell.text = summary_data[r][c]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(13)
            p.font.name = "Calibri"
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
            else:
                p.font.color.rgb = DARK
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY if r == 0 else WHITE

# Method comparison
add_shape(slide, Inches(0.8), Inches(4.3), Inches(5.5), Inches(2.8), WHITE)
add_text_box(slide, Inches(1.2), Inches(4.5), Inches(4.8), Inches(0.4),
             "方法学对比", 18, TEAL, bold=True)
meth_items = [
    "α,β-CROWN：线性松弛 + BaB",
    "  → 速度跨ε稳定 (~0.4s)",
    "  → Incomplete bounds 在中等难度",
    "     可被 Marabou complete 超越",
    "",
    "Marabou：SMT/SAT 求解",
    "  → Complete (sound+complete)",
    "  → ε=0.03 方差大 (3.5s ~ 931s)",
    "  → 速度 8.8~543× 慢于 BaB",
]
tb_m = add_text_box(slide, Inches(1.2), Inches(5.0), Inches(4.8), Inches(1.8), "", 12, DARK)
tf_m = tb_m.text_frame
tf_m.word_wrap = True
for item in meth_items:
    add_para(tf_m, item, 12, DARK if not item.startswith("  ") else GRAY, space_before=1)

# Value box
add_shape(slide, Inches(7.0), Inches(4.3), Inches(5.5), Inches(2.8), WHITE)
add_text_box(slide, Inches(7.4), Inches(4.5), Inches(4.8), Inches(0.4),
             "跨工具验证的价值", 18, GREEN, bold=True)
val_items = [
    "两条独立技术路线相互确认：",
    "  低 ε 完美一致 → 结论可信",
    "  分歧揭示 incomplete 边界",
    "  共同超时标示能力边界",
    "",
    "CIFAR-10 双重确认：2 个工具均",
    "触达 safe 证明边界（ε=2/255）",
]
tb_v = add_text_box(slide, Inches(7.4), Inches(5.0), Inches(4.8), Inches(1.8), "", 12, DARK)
tf_v = tb_v.text_frame
tf_v.word_wrap = True
for item in val_items:
    add_para(tf_v, item, 12, DARK if not item.startswith("  ") else GRAY, space_before=1)

slide_number(slide, 10, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════
# SLIDE 11: Paper Reproduction
# ══════════════════════════════════════════════════════
slide = make_slide()
add_bg(slide, OFF_WHITE)
add_shape(slide, 0, 0, W, Inches(1.0), NAVY)
add_text_box(slide, Inches(0.8), Inches(0.15), Inches(10), Inches(0.7),
             "论文实验对照：5/5 核心结论复现", 28, WHITE, bold=True)

papers = [
    ("α,β-CROWN (Xu et al., ICLR'21)", "CIFAR-10 CNN 验证速度 + 消融", TEAL),
    ("VNN-COMP 2021 竞赛报告", "mnistfc benchmark 排名", TEAL),
]
for i, (name, desc, color) in enumerate(papers):
    x = Inches(0.8) + Inches(i * 6.0)
    add_shape(slide, x, Inches(1.3), Inches(5.8), Inches(0.6), color)
    add_text_box(slide, x + Inches(0.2), Inches(1.35), Inches(5.4), Inches(0.3),
                 name, 14, WHITE, bold=True)

# Five core conclusions
conclusions = [
    ("① CROWN 速度 ~0.2s 且不随 ε 增长", "M6: 0.20–0.26s 恒定于 4 个 ε", "✅"),
    ("② α-CROWN 边界比 CROWN 更紧", "M6: VRA +1~5pp（ε 相关）", "✅"),
    ("③ BaB complete >> incomplete", "M4: paper config 80% vs 46% (ε=0.03)", "✅"),
    ("④ α,β-CROWN > Marabou (排名一致)", "M8: 4/5 vs 3/5, 0.66s vs 45.7s", "✅"),
    ("⑤ CIFAR-10 >> MNIST 难度", "M7+M8: 98% → 0% VRA 断层", "✅"),
]
for i, (claim, evidence, status) in enumerate(conclusions):
    y = Inches(2.2) + Inches(i * 0.85)
    add_shape(slide, Inches(0.8), y, Inches(11.5), Inches(0.7), WHITE)
    add_shape(slide, Inches(0.8), y, Inches(0.1), Inches(0.7), GREEN if status == "✅" else AMBER)
    add_text_box(slide, Inches(1.2), y + Inches(0.05), Inches(5.5), Inches(0.4),
                 claim, 14, DARK, bold=True)
    add_text_box(slide, Inches(1.2), y + Inches(0.4), Inches(7), Inches(0.3),
                 evidence, 12, GRAY)
    add_text_box(slide, Inches(11.0), y + Inches(0.1), Inches(1), Inches(0.5),
                 status, 20, GREEN, bold=True, alignment=PP_ALIGN.CENTER)

# Footnote
add_shape(slide, Inches(0.8), Inches(6.6), Inches(11.5), Inches(0.5), RGBColor(0xE8, 0xF0, 0xF8))
add_text_box(slide, Inches(1.2), Inches(6.65), Inches(10.8), Inches(0.4),
             "定性一致已完整验证；1:1 数值精确复现需要相同 benchmark 实例 + MIP 精化，已在 §4.3.5 诚实声明差异",
             11, GRAY)

slide_number(slide, 11, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════
# SLIDE 12: Paper Config Benchmark
# ══════════════════════════════════════════════════════
slide = make_slide()
add_bg(slide, OFF_WHITE)
add_shape(slide, 0, 0, W, Inches(1.0), NAVY)
add_text_box(slide, Inches(0.8), Inches(0.15), Inches(10), Inches(0.7),
             "论文配置对标实验", 28, WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.3), Inches(6), Inches(0.4),
             "BaB 参数：β-CROWN iter=20, lr_beta=0.03, batch=4096, kfsb_c5, timeout=120s", 14, TEAL)

paper_bench = [
    ["ε", "Paper config VRA", "耗时 (safe)", "timeout", "论文对照"],
    ["0.02", "5/5 (100%)", "0.24s", "0", "全部 BaB 直接安全"],
    ["0.03", "4/5 (80%)", "0.66s", "1 (sample 2)", "1 样本超时"],
]
r7, c7 = len(paper_bench), len(paper_bench[0])
tbl7 = slide.shapes.add_table(r7, c7, Inches(0.8), Inches(2.0), Inches(7.5), Inches(1.5)).table
for r in range(r7):
    for c in range(c7):
        cell = tbl7.cell(r, c)
        cell.text = paper_bench[r][c]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(13)
            p.font.name = "Calibri"
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
            else:
                p.font.color.rgb = DARK
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY if r == 0 else WHITE

# vs Marabou
add_text_box(slide, Inches(0.8), Inches(3.8), Inches(6), Inches(0.4),
             "同一 5 样本，ε=0.03：α,β-CROWN vs Marabou", 16, NAVY, bold=True)

compare_bench = [
    ["工具", "VRA", "平均耗时 (safe)", "timeout"],
    ["α,β-CROWN BaB (paper)", "4/5", "0.66s", "1"],
    ["Marabou SMT", "3/5", "45.7s", "2"],
]
r8, c8 = len(compare_bench), len(compare_bench[0])
tbl8 = slide.shapes.add_table(r8, c8, Inches(0.8), Inches(4.3), Inches(7.5), Inches(1.3)).table
for r in range(r8):
    for c in range(c8):
        cell = tbl8.cell(r, c)
        cell.text = compare_bench[r][c]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(13)
            p.font.name = "Calibri"
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
            elif r == 1:
                p.font.bold = True
                p.font.color.rgb = GREEN
            else:
                p.font.color.rgb = DARK
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY if r == 0 else WHITE

# Value proposition
add_shape(slide, Inches(9.0), Inches(1.3), Inches(3.8), Inches(5.8), NAVY)
add_text_box(slide, Inches(9.4), Inches(1.5), Inches(3.0), Inches(0.4),
             "对标意义", 20, WHITE, bold=True)
tb_bench = add_text_box(slide, Inches(9.4), Inches(2.1), Inches(3.0), Inches(4.5), "", 13, OFF_WHITE)
tf_bench = tb_bench.text_frame
tf_bench.word_wrap = True
add_para(tf_bench, "方向完全一致：", 14, AMBER, bold=True, space_before=4)
add_para(tf_bench, "α,β-CROWN > Marabou", 13, WHITE, bold=True, space_before=2)
add_para(tf_bench, "与 VNN-COMP 2021", 13, OFF_WHITE)
add_para(tf_bench, "Table 16 排名一致", 13, OFF_WHITE)
add_para(tf_bench, "", 8, OFF_WHITE)
add_para(tf_bench, "速度优势 69×:", 14, AMBER, bold=True, space_before=4)
add_para(tf_bench, "0.66s vs 45.7s", 13, WHITE, bold=True)
add_para(tf_bench, "bound propagation", 13, OFF_WHITE)
add_para(tf_bench, "优于 SMT 求解", 13, OFF_WHITE)
add_para(tf_bench, "", 8, OFF_WHITE)
add_para(tf_bench, "消融验证：", 14, AMBER, bold=True, space_before=4)
add_para(tf_bench, "iter=20 vs 50 无差异", 13, OFF_WHITE)
add_para(tf_bench, "batch_size不影响", 13, OFF_WHITE)
add_para(tf_bench, "简单实例", 13, OFF_WHITE)

slide_number(slide, 12, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════
# SLIDE 13: Capability Spectrum
# ══════════════════════════════════════════════════════
slide = make_slide()
add_bg(slide, OFF_WHITE)
add_shape(slide, 0, 0, W, Inches(1.0), NAVY)
add_text_box(slide, Inches(0.8), Inches(0.15), Inches(10), Inches(0.7),
             "验证策略能力光谱", 28, WHITE, bold=True)

spec_data = [
    ["方法", "算法类型", "ε=0.03 VRA", "耗时", "证明 safe", "找 unsafe"],
    ["CROWN", "固定线性松弛", "41%", "0.26s", "✅", "❌"],
    ["α-CROWN", "可优化松弛", "46%", "1.58s", "✅", "❌"],
    ["PGD", "梯度攻击", "—", "~1s", "❌", "✅"],
    ["BaB (babsr)", "分支定界(基线)", "65%", "14.27s", "✅", "✅"],
    ["BaB (kfsb_c5)", "分支定界(优化)", "68%", "6.66s", "✅", "✅"],
    ["PGD+BaB(kfsb)", "PGD预筛+BaB", "73%", "6.45s", "✅", "✅"],
    ["Marabou SMT", "SMT求解", "3/5*", "45.7s", "✅", "✅"],
]
r9, c9 = len(spec_data), len(spec_data[0])
tbl9 = slide.shapes.add_table(r9, c9, Inches(0.5), Inches(1.3), Inches(12.3), Inches(4.0)).table
for r in range(r9):
    for c in range(c9):
        cell = tbl9.cell(r, c)
        cell.text = spec_data[r][c]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(12)
            p.font.name = "Calibri"
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
            elif r == 6:
                p.font.bold = True
                p.font.color.rgb = NAVY
            else:
                p.font.color.rgb = DARK
        cell.fill.solid()
        if r == 0:
            cell.fill.fore_color.rgb = NAVY
        elif r == 6:
            cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xF8)
        else:
            cell.fill.fore_color.rgb = WHITE

add_shape(slide, Inches(0.5), Inches(5.6), Inches(12.3), Inches(1.5), WHITE)
add_text_box(slide, Inches(1.0), Inches(5.7), Inches(11), Inches(0.4),
             "按 ε 分区的工程推荐方案", 18, TEAL, bold=True)

recs = [
    ("≤0.01", "CROWN 快速筛查", "0.20s, 98% VRA，无需 BaB"),
    ("0.02~0.03", "PGD + BaB (kfsb_c5)", "VRA 最高, timeout 最低"),
    ("≥0.05", "PGD + BaB (kfsb_c5)", "可执行性优先"),
    ("独立复核", "Marabou SMT", "完备 SMT 交叉验证"),
]
for i, (eps, strat, reason) in enumerate(recs):
    x = Inches(0.6) + Inches(i * 3.15)
    add_text_box(slide, x, Inches(6.2), Inches(2.8), Inches(0.3),
                 f"ε {eps}", 12, AMBER, bold=True)
    add_text_box(slide, x, Inches(6.45), Inches(2.8), Inches(0.3),
                 strat, 12, DARK, bold=True)
    add_text_box(slide, x, Inches(6.75), Inches(2.8), Inches(0.3),
                 reason, 10, GRAY)

slide_number(slide, 13, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════
# SLIDE 14: Key Conclusions 1
# ══════════════════════════════════════════════════════
slide = make_slide()
add_bg(slide, OFF_WHITE)
add_shape(slide, 0, 0, W, Inches(1.0), NAVY)
add_text_box(slide, Inches(0.8), Inches(0.15), Inches(10), Inches(0.7),
             "核心结论（1/2）", 28, WHITE, bold=True)

conclusions_left = [
    ("结论 1: ADO 方法论有效性", TEAL,
     "受控消融-量化分解-最优配置三步骤在 method×candidates×reduceop 上实现 BaB 系统性优化",
     "VRA +2.0pp, timeout −2, Mean Time −0.81s"),
    ("结论 2: PGD + BaB 互补关系", TEAL,
     "PGD unsafe 召回率 100%（覆盖 BaB 全部反例），预筛使 timeout 降 62~69%",
     "ε=0.03 公平对照 VRA 从 68%→73%"),
    ("结论 3: 不完整验证能力边界", TEAL,
     "CROWN ~0.2s 恒定与 ε 无关；α-CROWN 倒U型收益在 ε=0.03 达峰值 +5pp",
     "ε≤0.01 用 CROWN 即可, ε=0.02~0.03 用 PGD+BaB"),
]
for i, (title, color, main, metric) in enumerate(conclusions_left):
    y = Inches(1.3) + Inches(i * 2.0)
    add_shape(slide, Inches(0.8), y, Inches(11.5), Inches(1.7), WHITE)
    add_shape(slide, Inches(0.8), y, Inches(0.1), Inches(1.7), color)
    add_text_box(slide, Inches(1.3), y + Inches(0.1), Inches(10.5), Inches(0.4),
                 title, 18, color, bold=True)
    tb = add_text_box(slide, Inches(1.3), y + Inches(0.6), Inches(7.5), Inches(0.6),
                     main, 13, DARK)
    tb.text_frame.word_wrap = True
    add_text_box(slide, Inches(9.0), y + Inches(0.6), Inches(3.0), Inches(0.8),
                 metric, 14, NAVY, bold=True, alignment=PP_ALIGN.RIGHT)

slide_number(slide, 14, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════
# SLIDE 15: Key Conclusions 2
# ══════════════════════════════════════════════════════
slide = make_slide()
add_bg(slide, OFF_WHITE)
add_shape(slide, 0, 0, W, Inches(1.0), NAVY)
add_text_box(slide, Inches(0.8), Inches(0.15), Inches(10), Inches(0.7),
             "核心结论（2/2）", 28, WHITE, bold=True)

conclusions_right = [
    ("结论 4: 跨ε稳定性", TEAL,
     "kfsb 优势呈倒U型，ε=0.03 达峰值 (VRA +3pp, 时间仅 baseline 47%)",
     "ε≥0.05 超时主导"),
    ("结论 5: CIFAR-10 触及能力边界", TEAL,
     "5 因素复合作用使 VRA 从 98% 断层至 0%；PGD 仍有效 (80% unsafe 检出)",
     "safe 证明功能已到边界"),
    ("结论 6: Marabou 跨工具确认", TEAL,
     "3 ε × 5 样本系统对照：低 ε 完美一致，中 ε 工具分歧，高 ε 共同边界",
     "α,β-CROWN 速度优势稳定"),
]
for i, (title, color, main, metric) in enumerate(conclusions_right):
    y = Inches(1.3) + Inches(i * 2.0)
    add_shape(slide, Inches(0.8), y, Inches(11.5), Inches(1.7), WHITE)
    add_shape(slide, Inches(0.8), y, Inches(0.1), Inches(1.7), color)
    add_text_box(slide, Inches(1.3), y + Inches(0.1), Inches(10.5), Inches(0.4),
                 title, 18, color, bold=True)
    tb = add_text_box(slide, Inches(1.3), y + Inches(0.6), Inches(7.5), Inches(0.6),
                     main, 13, DARK)
    tb.text_frame.word_wrap = True
    add_text_box(slide, Inches(9.0), y + Inches(0.6), Inches(3.0), Inches(0.8),
                 metric, 14, NAVY, bold=True, alignment=PP_ALIGN.RIGHT)

slide_number(slide, 15, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════
# SLIDE 16: Limitations
# ══════════════════════════════════════════════════════
slide = make_slide()
add_bg(slide, OFF_WHITE)
add_shape(slide, 0, 0, W, Inches(1.0), NAVY)
add_text_box(slide, Inches(0.8), Inches(0.15), Inches(10), Inches(0.7),
             "当前局限", 28, WHITE, bold=True)

limits = [
    ("CIFAR-10 最小可行实验", "24 样本、2 ε、1 完整配置 — 未做分支消融和 ε 网格"),
    ("Marabou 5 样本规模", "发现分歧现象但不足以做统计显著的量化结论"),
    ("ADO 单模型验证", "MNIST FCNN 单一架构，跨模型泛化性待检验"),
    ("消融维度有限", "仅覆盖 method × candidates × reduceop，其他参数未纳入"),
    ("无 MIP 精化", "无 Gurobi license，无法使用 bab-refine 增强"),
]
for i, (title, desc) in enumerate(limits):
    y = Inches(1.3) + Inches(i * 1.15)
    add_shape(slide, Inches(0.8), y, Inches(11.5), Inches(0.95), WHITE)
    add_shape(slide, Inches(0.8), y, Inches(0.08), Inches(0.95), AMBER)
    add_text_box(slide, Inches(1.2), y + Inches(0.1), Inches(10.5), Inches(0.35),
                 title, 16, DARK, bold=True)
    add_text_box(slide, Inches(1.2), y + Inches(0.5), Inches(10.5), Inches(0.35),
                 desc, 12, GRAY)

# Honesty statement
add_shape(slide, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.5), NAVY)
add_text_box(slide, Inches(1.2), Inches(6.15), Inches(10.5), Inches(0.4),
             "所有结论标注了证据来源 → §7 证据链索引可逐条溯源 → 数据文件提交时一并附上",
             12, OFF_WHITE)

slide_number(slide, 16, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════
# SLIDE 17: Experiment Summary (by numbers)
# ══════════════════════════════════════════════════════
slide = make_slide()
add_bg(slide, NAVY)
add_text_box(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6),
             "实验规模汇总", 28, WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.3),
             "46 组配置 · 2 数据集 · 2 独立验证工具 · 2 篇对标论文", 14, RGBColor(0xBB, 0xCC, 0xDD))

big_nums = [
    ("46", "配置组数", "覆盖 6 种策略 × 4 ε"),
    ("15", "跨工具对照", "5 样本 × 3 ε × 2 工具"),
    ("2", "数据集", "MNIST + CIFAR-10"),
    ("5/5", "论文结论", "全部定性复现"),
]
for i, (num, label, sub) in enumerate(big_nums):
    x = Inches(0.8) + Inches(i * 3.2)
    add_shape(slide, x, Inches(1.6), Inches(2.8), Inches(2.2), RGBColor(0x15, 0x30, 0x4F))
    add_text_box(slide, x, Inches(1.8), Inches(2.8), Inches(1.0),
                 num, 44, AMBER, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(2.8), Inches(2.8), Inches(0.4),
                 label, 16, WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(3.2), Inches(2.8), Inches(0.3),
                 sub, 11, RGBColor(0xBB, 0xCC, 0xDD), alignment=PP_ALIGN.CENTER)

# Evidence chain summary
add_shape(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(2.8), RGBColor(0x15, 0x30, 0x4F))
add_text_box(slide, Inches(1.2), Inches(4.4), Inches(5), Inches(0.4),
             "证据链索引关键条目", 18, AMBER, bold=True)

evidences = [
    "不完整验证: m6_incomplete_compare.csv (8 组)",
    "PGD+BaB: m5_pgd_compare.csv + m5_pgd_control.csv",
    "分支消融: m3_branching_ablation.csv (5 组)",
    "ε 网格: m4_epsilon_grid.csv (12 组)",
    "跨工具对照: m8_marabou/*.csv (ε=0.01/0.02/0.03)",
    "论文对标: exp_configs/course/m8_paper_repro/*.yaml",
]
for i, item in enumerate(evidences):
    y = Inches(4.9) + Inches(i * 0.35)
    add_text_box(slide, Inches(1.2), y, Inches(10.5), Inches(0.3),
                 f"  >  {item}", 11, OFF_WHITE)

slide_number(slide, 17, TOTAL_SLIDES)

# ══════════════════════════════════════════════════════
# SLIDE 18: Thank You
# ══════════════════════════════════════════════════════
slide = make_slide()
add_bg(slide, NAVY)
add_shape(slide, 0, Inches(2.5), W, Inches(2.5), TEAL)
add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
             "谢谢！", 48, WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.5),
             "欢迎提问与讨论", 20, OFF_WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.4),
             "α,β-CROWN + Marabou · MNIST + CIFAR-10 · ADO 方法论 · 6 策略光谱",
             12, RGBColor(0xBB, 0xCC, 0xDD), alignment=PP_ALIGN.CENTER)
slide_number(slide, 18, TOTAL_SLIDES)

# ── Save ──
out_path = os.path.expanduser("~/alpha-beta-CROWN/项目书/results/完成情况汇报.pptx")
prs.save(out_path)
print(f"PPT saved to: {out_path}")
print(f"Total slides: {len(prs.slides)}")
